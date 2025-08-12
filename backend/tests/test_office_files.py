"""
Officeファイル対応テスト
実行日時: 2025-08-13

テスト内容:
- DOCX/XLSX/PPTX/PDF/TXTファイルのアップロード
- テキスト抽出機能
"""

import asyncio
import aiohttp
import json
import io
from datetime import datetime
from typing import Dict, Any

# テスト用ライブラリ
import docx
import openpyxl
from pptx import Presentation
import PyPDF2

BASE_URL = "http://localhost:8000"


class OfficeFilesTest:
    """Officeファイル対応のテストクラス"""
    
    def __init__(self):
        self.session = None
        self.test_results = []
        self.library_id = None
        self.created_files = []
        
    async def setup(self):
        """セットアップ"""
        self.session = aiohttp.ClientSession()
        
    async def teardown(self):
        """クリーンアップ"""
        # 作成したファイルを削除
        if self.library_id:
            for filename in self.created_files:
                try:
                    await self.session.delete(
                        f"{BASE_URL}/api/libraries/{self.library_id}/files/{filename}"
                    )
                except:
                    pass
            
            # ライブラリを削除
            try:
                await self.session.delete(f"{BASE_URL}/api/libraries/{self.library_id}")
            except:
                pass
                
        if self.session:
            await self.session.close()
    
    def create_test_docx(self) -> bytes:
        """テスト用DOCXファイルを作成"""
        doc = docx.Document()
        doc.add_heading('テストドキュメント', 0)
        doc.add_paragraph('これはテスト用のWordドキュメントです。')
        doc.add_heading('セクション1', level=1)
        doc.add_paragraph('日本語のテキストも正しく抽出されることを確認します。')
        
        # テーブルを追加
        table = doc.add_table(rows=3, cols=3)
        table.cell(0, 0).text = 'ヘッダー1'
        table.cell(0, 1).text = 'ヘッダー2'
        table.cell(0, 2).text = 'ヘッダー3'
        table.cell(1, 0).text = 'データ1'
        table.cell(1, 1).text = 'データ2'
        table.cell(1, 2).text = 'データ3'
        
        # バイト列として保存
        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()
    
    def create_test_xlsx(self) -> bytes:
        """テスト用XLSXファイルを作成"""
        workbook = openpyxl.Workbook()
        
        # シート1
        sheet1 = workbook.active
        sheet1.title = "売上データ"
        sheet1['A1'] = '商品名'
        sheet1['B1'] = '数量'
        sheet1['C1'] = '価格'
        sheet1['A2'] = 'りんご'
        sheet1['B2'] = 100
        sheet1['C2'] = 150
        sheet1['A3'] = 'みかん'
        sheet1['B3'] = 200
        sheet1['C3'] = 100
        
        # シート2
        sheet2 = workbook.create_sheet("在庫データ")
        sheet2['A1'] = '商品コード'
        sheet2['B1'] = '在庫数'
        sheet2['A2'] = 'P001'
        sheet2['B2'] = 500
        
        # バイト列として保存
        buffer = io.BytesIO()
        workbook.save(buffer)
        return buffer.getvalue()
    
    def create_test_pptx(self) -> bytes:
        """テスト用PPTXファイルを作成"""
        prs = Presentation()
        
        # スライド1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "プレゼンテーションタイトル"
        slide1.placeholders[1].text = "サブタイトル"
        
        # スライド2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "アジェンダ"
        content = slide2.placeholders[1]
        content.text = "1. はじめに\n2. 本編\n3. まとめ"
        
        # バイト列として保存
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
    
    def create_test_pdf(self) -> bytes:
        """簡単なテスト用PDFを作成（テキストベース）"""
        # 簡単なPDFテキストファイルを作成
        pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
/Resources << /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >>
/Contents 4 0 R >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT
/F1 12 Tf
100 700 Td
(Test PDF Document) Tj
ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000274 00000 n
trailer
<< /Size 5 /Root 1 0 R >>
startxref
365
%%EOF"""
        return pdf_content
    
    async def test_create_library(self) -> Dict[str, Any]:
        """テスト用ライブラリを作成"""
        test_name = "Create Test Library"
        try:
            data = {
                "name": f"Officeファイルテスト_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                "description": "Officeファイル対応のテスト用ライブラリ"
            }
            
            async with self.session.post(
                f"{BASE_URL}/api/libraries",
                json=data
            ) as response:
                result = await response.json()
                
                if response.status == 200:
                    self.library_id = result['id']
                    return {
                        "test": test_name,
                        "status": "success",
                        "library_id": self.library_id
                    }
                else:
                    return {
                        "test": test_name,
                        "status": "failed",
                        "error": f"Status {response.status}"
                    }
                    
        except Exception as e:
            return {
                "test": test_name,
                "status": "error",
                "error": str(e)
            }
    
    async def test_upload_and_extract(self, file_type: str, content: bytes) -> Dict[str, Any]:
        """ファイルアップロードとテキスト抽出テスト"""
        test_name = f"Upload and Extract {file_type.upper()}"
        filename = f"test.{file_type}"
        
        try:
            # ファイルアップロード
            data = aiohttp.FormData()
            data.add_field('file', content, filename=filename, 
                          content_type='application/octet-stream')
            
            async with self.session.post(
                f"{BASE_URL}/api/libraries/{self.library_id}/files",
                data=data
            ) as response:
                if response.status != 200:
                    return {
                        "test": test_name,
                        "status": "failed",
                        "stage": "upload",
                        "error": f"Upload failed: {response.status}"
                    }
                
                self.created_files.append(filename)
            
            # テキスト抽出
            async with self.session.get(
                f"{BASE_URL}/api/libraries/{self.library_id}/files/{filename}/text"
            ) as response:
                result = await response.json()
                
                if response.status == 200:
                    text_length = len(result.get('text', ''))
                    metadata = result.get('metadata', {})
                    
                    return {
                        "test": test_name,
                        "status": "success",
                        "text_length": text_length,
                        "format": metadata.get('format'),
                        "metadata": metadata,
                        "sample_text": result.get('text', '')[:200]  # 最初の200文字
                    }
                else:
                    return {
                        "test": test_name,
                        "status": "failed",
                        "stage": "extract",
                        "error": f"Extract failed: {response.status}",
                        "detail": result
                    }
                    
        except Exception as e:
            return {
                "test": test_name,
                "status": "error",
                "error": str(e)
            }
    
    async def run_all_tests(self):
        """全テストを実行"""
        print("\n" + "="*60)
        print("Officeファイル対応テスト開始")
        print("="*60)
        
        await self.setup()
        
        try:
            # 1. ライブラリ作成
            print("\n[1] テスト用ライブラリ作成...")
            result = await self.test_create_library()
            self.test_results.append(result)
            print(f"   結果: {result['status']}")
            
            if result['status'] != 'success':
                print("   ライブラリ作成失敗のため、テストを中止")
                return
            
            # 2. TXTファイルテスト
            print("\n[2] TXTファイルテスト...")
            txt_content = "これはテスト用のテキストファイルです。\n日本語の処理を確認します。".encode('utf-8')
            result = await self.test_upload_and_extract('txt', txt_content)
            self.test_results.append(result)
            print(f"   結果: {result['status']}")
            if result['status'] == 'success':
                print(f"   抽出文字数: {result['text_length']}")
            
            # 3. DOCXファイルテスト
            print("\n[3] DOCXファイルテスト...")
            docx_content = self.create_test_docx()
            result = await self.test_upload_and_extract('docx', docx_content)
            self.test_results.append(result)
            print(f"   結果: {result['status']}")
            if result['status'] == 'success':
                print(f"   抽出文字数: {result['text_length']}")
                print(f"   段落数: {result['metadata'].get('paragraph_count')}")
                print(f"   テーブル数: {result['metadata'].get('table_count')}")
            
            # 4. XLSXファイルテスト
            print("\n[4] XLSXファイルテスト...")
            xlsx_content = self.create_test_xlsx()
            result = await self.test_upload_and_extract('xlsx', xlsx_content)
            self.test_results.append(result)
            print(f"   結果: {result['status']}")
            if result['status'] == 'success':
                print(f"   抽出文字数: {result['text_length']}")
                print(f"   シート数: {result['metadata'].get('sheet_count')}")
                print(f"   シート名: {result['metadata'].get('sheet_names')}")
            
            # 5. PPTXファイルテスト
            print("\n[5] PPTXファイルテスト...")
            pptx_content = self.create_test_pptx()
            result = await self.test_upload_and_extract('pptx', pptx_content)
            self.test_results.append(result)
            print(f"   結果: {result['status']}")
            if result['status'] == 'success':
                print(f"   抽出文字数: {result['text_length']}")
                print(f"   スライド数: {result['metadata'].get('slide_count')}")
            
            # 6. PDFファイルテスト
            print("\n[6] PDFファイルテスト...")
            pdf_content = self.create_test_pdf()
            result = await self.test_upload_and_extract('pdf', pdf_content)
            self.test_results.append(result)
            print(f"   結果: {result['status']}")
            if result['status'] == 'success':
                print(f"   抽出文字数: {result['text_length']}")
                print(f"   ページ数: {result['metadata'].get('page_count')}")
                
        finally:
            await self.teardown()
        
        # 結果サマリー
        self.print_summary()
        
        # 結果をJSONファイルに保存
        self.save_results()
    
    def print_summary(self):
        """テスト結果サマリーを表示"""
        print("\n" + "="*60)
        print("テスト結果サマリー")
        print("="*60)
        
        success_count = sum(1 for r in self.test_results if r['status'] == 'success')
        failed_count = sum(1 for r in self.test_results if r['status'] == 'failed')
        error_count = sum(1 for r in self.test_results if r['status'] == 'error')
        total_count = len(self.test_results)
        
        print(f"総テスト数: {total_count}")
        print(f"成功: {success_count}")
        print(f"失敗: {failed_count}")
        print(f"エラー: {error_count}")
        
        # 対応フォーマット
        print("\n【対応確認済みフォーマット】")
        for result in self.test_results:
            if result['status'] == 'success' and 'format' in result:
                print(f"  ✅ {result['format']}")
        
        if failed_count > 0 or error_count > 0:
            print("\n【失敗/エラーのテスト】")
            for result in self.test_results:
                if result['status'] in ['failed', 'error']:
                    print(f"  - {result['test']}: {result.get('error', 'Unknown error')}")
    
    def save_results(self):
        """テスト結果をJSONファイルに保存"""
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"test_office_files_results_{timestamp}.json"
        
        with open(filename, 'w', encoding='utf-8') as f:
            json.dump({
                'test_date': datetime.now().isoformat(),
                'results': self.test_results
            }, f, ensure_ascii=False, indent=2)
        
        print(f"\nテスト結果を {filename} に保存しました")


async def main():
    """メイン実行関数"""
    tester = OfficeFilesTest()
    await tester.run_all_tests()


if __name__ == "__main__":
    asyncio.run(main())