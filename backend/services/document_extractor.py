"""
ドキュメント抽出サービス
各種ファイル形式からテキストを抽出する

対応形式:
- PDF
- TXT
- DOCX (Microsoft Word)
- XLSX (Microsoft Excel)  
- PPTX (Microsoft PowerPoint)
"""

import os
import io
from typing import Optional, List, Union
from dataclasses import dataclass
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
from utils.logger import api_logger


@dataclass
class DocumentMetadata:
    """ドキュメントのメタデータ"""
    format: str
    page_count: Optional[int] = None
    paragraph_count: Optional[int] = None
    table_count: Optional[int] = None
    sheet_count: Optional[int] = None
    sheet_names: Optional[List[str]] = None
    slide_count: Optional[int] = None
    size: Optional[int] = None


@dataclass
class ExtractionResult:
    """テキスト抽出結果"""
    success: bool
    text: str
    metadata: DocumentMetadata
    error: Optional[str] = None


class DocumentExtractor:
    """各種ドキュメントからテキストを抽出するクラス"""
    
    @staticmethod
    async def extract_text(content: bytes, filename: str) -> ExtractionResult:
        """
        ファイルからテキストを抽出
        
        Args:
            content: ファイルのバイナリコンテンツ
            filename: ファイル名（拡張子判定用）
            
        Returns:
            抽出結果
        """
        # 拡張子を取得
        ext = os.path.splitext(filename)[1].lower()
        
        try:
            if ext == '.pdf':
                return await DocumentExtractor._extract_pdf(content)
            elif ext == '.txt':
                return await DocumentExtractor._extract_txt(content)
            elif ext == '.docx':
                return await DocumentExtractor._extract_docx(content)
            elif ext == '.xlsx':
                return await DocumentExtractor._extract_xlsx(content)
            elif ext == '.pptx':
                return await DocumentExtractor._extract_pptx(content)
            else:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata=DocumentMetadata(format='UNKNOWN'),
                    error=f'Unsupported file type: {ext}'
                )
        except Exception as e:
            api_logger.error(f"Error extracting text from {filename}: {str(e)}")
            return ExtractionResult(
                success=False,
                text='',
                metadata=DocumentMetadata(format=ext.upper()[1:] if ext else 'UNKNOWN'),
                error=str(e)
            )
    
    @staticmethod
    async def _extract_pdf(content: bytes) -> ExtractionResult:
        """PDFからテキストを抽出"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- ページ {page_num + 1} ---\n{page_text}")
            
            return ExtractionResult(
                success=True,
                text='\n\n'.join(text_parts),
                metadata=DocumentMetadata(
                    format='PDF',
                    page_count=len(pdf_reader.pages)
                )
            )
        except Exception as e:
            api_logger.error(f"PDF extraction error: {str(e)}")
            raise
    
    @staticmethod
    async def _extract_txt(content: bytes) -> ExtractionResult:
        """テキストファイルからテキストを抽出"""
        try:
            # 日本語対応のエンコーディングを優先的に試す
            encodings = ['utf-8', 'shift_jis', 'cp932', 'euc-jp', 'iso-2022-jp']
            text = None
            
            for encoding in encodings:
                try:
                    text = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                # フォールバック: エラーを無視してデコード
                text = content.decode('utf-8', errors='ignore')
            
            return ExtractionResult(
                success=True,
                text=text,
                metadata=DocumentMetadata(
                    format='TXT',
                    size=len(content)
                )
            )
        except Exception as e:
            api_logger.error(f"TXT extraction error: {str(e)}")
            raise
    
    @staticmethod
    async def _extract_docx(content: bytes) -> ExtractionResult:
        """DOCXからテキストを抽出"""
        try:
            docx_file = io.BytesIO(content)
            doc = docx.Document(docx_file)
            
            text_parts = []
            
            # 段落を抽出
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # テーブルを抽出
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    if any(row_text):
                        table_text.append(' | '.join(row_text))
                if table_text:
                    text_parts.append('\n'.join(table_text))
            
            return ExtractionResult(
                success=True,
                text='\n\n'.join(text_parts),
                metadata=DocumentMetadata(
                    format='DOCX',
                    paragraph_count=len(doc.paragraphs),
                    table_count=len(doc.tables)
                )
            )
        except Exception as e:
            api_logger.error(f"DOCX extraction error: {str(e)}")
            raise
    
    @staticmethod
    async def _extract_xlsx(content: bytes) -> ExtractionResult:
        """XLSXからテキストを抽出"""
        try:
            xlsx_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                sheet_data.append(f"=== シート: {sheet_name} ===")
                
                # 各行のデータを抽出
                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        if cell.value is not None:
                            row_data.append(str(cell.value))
                    if row_data:
                        sheet_data.append(' | '.join(row_data))
                
                if len(sheet_data) > 1:  # ヘッダー以外にデータがある場合
                    text_parts.append('\n'.join(sheet_data))
            
            return ExtractionResult(
                success=True,
                text='\n\n'.join(text_parts),
                metadata=DocumentMetadata(
                    format='XLSX',
                    sheet_count=len(workbook.sheetnames),
                    sheet_names=list(workbook.sheetnames)
                )
            )
        except Exception as e:
            api_logger.error(f"XLSX extraction error: {str(e)}")
            raise
    
    @staticmethod
    async def _extract_pptx(content: bytes) -> ExtractionResult:
        """PPTXからテキストを抽出"""
        try:
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"--- スライド {slide_num} ---"]
                
                # タイトル
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        slide_text.append(f"タイトル: {title}")
                
                # 各シェイプからテキストを抽出
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        # タイトルの重複を避ける
                        if text and (not slide.shapes.title or text != slide.shapes.title.text):
                            slide_text.append(text)
                    
                    # テーブルの場合
                    if shape.has_table:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(' | '.join(row_text))
                        if table_text:
                            slide_text.append('\n'.join(table_text))
                
                if len(slide_text) > 1:  # ヘッダー以外にコンテンツがある場合
                    text_parts.append('\n'.join(slide_text))
            
            return ExtractionResult(
                success=True,
                text='\n\n'.join(text_parts),
                metadata=DocumentMetadata(
                    format='PPTX',
                    slide_count=len(presentation.slides)
                )
            )
        except Exception as e:
            api_logger.error(f"PPTX extraction error: {str(e)}")
            raise


# シングルトンインスタンス
document_extractor = DocumentExtractor()